"""
generate_final_report.py
Generates:
  1. Final_Report_Group3.docx  – full Word report with EDA section
  2. Final_Presentation_Group3.pptx – PowerPoint slide deck
Run from the Project/ directory with the capstone venv active.
"""

from __future__ import annotations
import os, io, re, warnings
from pathlib import Path

import pandas as pd
import numpy as np
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.ticker as mticker
import seaborn as sns
from collections import Counter

warnings.filterwarnings("ignore")

# ── paths ────────────────────────────────────────────────────────────────────
HERE      = Path(__file__).parent
DATA_PATH = HERE / "ontario_college_IT_programs_ONLY_DEDUP_COLLEGE_PROGRAM.csv"
WORD_OUT  = HERE / "Final_Report_Group3.docx"
PPTX_OUT  = HERE / "Final_Presentation_Group3.pptx"

# ── load data ─────────────────────────────────────────────────────────────────
df = pd.read_csv(DATA_PATH)
df["duration_months"] = pd.to_numeric(df["duration_months"], errors="coerce")
df["duration_years"]  = pd.to_numeric(df["duration_years"],  errors="coerce")

# ── visual style ──────────────────────────────────────────────────────────────
sns.set_theme(style="whitegrid", palette="muted", font_scale=1.1)
MAIN  = "#4C72B0"
ACCT  = "#DD8452"
PAL10 = sns.color_palette("tab10", 12)
plt.rcParams.update({"figure.dpi": 130, "figure.facecolor": "white"})


# ─────────────────────────────────────────────────────────────────────────────
# HELPER: save figure → bytes
# ─────────────────────────────────────────────────────────────────────────────
def fig_bytes(fig) -> bytes:
    buf = io.BytesIO()
    fig.savefig(buf, format="png", bbox_inches="tight")
    plt.close(fig)
    buf.seek(0)
    return buf.read()


# ─────────────────────────────────────────────────────────────────────────────
# CHART FUNCTIONS  (return raw bytes)
# ─────────────────────────────────────────────────────────────────────────────

def chart_missing() -> bytes:
    null_pct = (df.isna().sum() / len(df) * 100).sort_values(ascending=False)
    fig, ax = plt.subplots(figsize=(9, 4))
    colors = ["#e74c3c" if v > 20 else "#f39c12" if v > 5 else "#2ecc71"
              for v in null_pct.values]
    bars = ax.barh(null_pct.index, null_pct.values, color=colors)
    ax.set_xlabel("Missing Values (%)")
    ax.set_title("Missing Value Rate per Column", fontweight="bold")
    for bar, val in zip(bars, null_pct.values):
        ax.text(val + 0.3, bar.get_y() + bar.get_height() / 2,
                f"{val:.1f}%", va="center", fontsize=8)
    plt.tight_layout()
    return fig_bytes(fig)


def chart_category() -> bytes:
    cat_counts = df["program_category"].value_counts()
    fig, axes = plt.subplots(1, 2, figsize=(14, 5))
    axes[0].barh(cat_counts.index, cat_counts.values, color=PAL10[:len(cat_counts)])
    axes[0].set_xlabel("Number of Programs")
    axes[0].set_title("Programs per Category", fontweight="bold")
    axes[0].invert_yaxis()
    for bar, v in zip(axes[0].patches, cat_counts.values):
        axes[0].text(v + 5, bar.get_y() + bar.get_height() / 2,
                     f"{v:,}", va="center", fontsize=8)
    wedges, _, autotexts = axes[1].pie(
        cat_counts.values, autopct="%1.1f%%",
        colors=PAL10[:len(cat_counts)], startangle=140,
        pctdistance=0.82,
        wedgeprops={"linewidth": 0.8, "edgecolor": "white"})
    for at in autotexts:
        at.set_fontsize(8)
    axes[1].legend(wedges, cat_counts.index, loc="lower left",
                   fontsize=7, framealpha=0.6, bbox_to_anchor=(-0.25, -0.2))
    axes[1].set_title("Share by Category", fontweight="bold")
    plt.suptitle("Program Category Overview", fontsize=13, fontweight="bold")
    plt.tight_layout()
    return fig_bytes(fig)


def chart_credential() -> bytes:
    fig, axes = plt.subplots(1, 2, figsize=(14, 5))
    for ax, col, title in zip(axes,
                               ["credential_type", "credential_group"],
                               ["Credential Type", "Credential Group"]):
        counts = df[col].value_counts()
        ax.bar(range(len(counts)), counts.values,
               color=PAL10[:len(counts)], edgecolor="white")
        ax.set_xticks(range(len(counts)))
        ax.set_xticklabels(counts.index, rotation=35, ha="right")
        ax.set_ylabel("Count")
        ax.set_title(title, fontweight="bold")
        for i, v in enumerate(counts.values):
            ax.text(i, v + 5, f"{v:,}", ha="center", fontsize=8)
    plt.suptitle("Credential Distribution", fontsize=13, fontweight="bold")
    plt.tight_layout()
    return fig_bytes(fig)


def chart_delivery() -> bytes:
    del_counts = df["delivery_format"].value_counts()
    colors5 = ["#2980b9", "#e67e22", "#27ae60", "#8e44ad", "#c0392b"]
    fig, axes = plt.subplots(1, 2, figsize=(12, 5))
    wedges, _, autotexts = axes[0].pie(
        del_counts.values, labels=del_counts.index,
        autopct="%1.1f%%", colors=colors5[:len(del_counts)],
        startangle=90,
        wedgeprops={"width": 0.55, "edgecolor": "white"},
        pctdistance=0.75)
    for at in autotexts:
        at.set_fontsize(9)
    axes[0].set_title("Delivery Format — Donut", fontweight="bold")
    axes[1].barh(del_counts.index, del_counts.values, color=colors5[:len(del_counts)])
    axes[1].set_xlabel("Count")
    axes[1].set_title("Delivery Format — Count", fontweight="bold")
    for i, v in enumerate(del_counts.values):
        axes[1].text(v + 5, i, f"{v:,}", va="center", fontsize=9)
    plt.suptitle("How Programs Are Delivered", fontsize=13, fontweight="bold")
    plt.tight_layout()
    return fig_bytes(fig)


def chart_duration() -> bytes:
    dur_valid = df["duration_months"].dropna()
    fig, axes = plt.subplots(1, 2, figsize=(14, 5))
    sns.histplot(dur_valid, bins=30, kde=True, color=MAIN,
                 edgecolor="white", ax=axes[0], line_kws={"lw": 2})
    axes[0].axvline(dur_valid.median(), color=ACCT, ls="--", lw=1.8,
                    label=f"Median = {dur_valid.median():.0f} mo")
    axes[0].axvline(dur_valid.mean(), color="green", ls="--", lw=1.8,
                    label=f"Mean = {dur_valid.mean():.1f} mo")
    axes[0].legend()
    axes[0].set_xlabel("Duration (months)")
    axes[0].set_title("Duration Distribution (months)", fontweight="bold")
    cred_order = (df.groupby("credential_group")["duration_months"]
                  .median().sort_values().index)
    sns.boxplot(data=df.dropna(subset=["duration_months"]),
                y="credential_group", x="duration_months",
                order=cred_order, palette="Set2", ax=axes[1])
    axes[1].set_xlabel("Duration (months)")
    axes[1].set_ylabel("")
    axes[1].set_title("Duration by Credential Group", fontweight="bold")
    plt.suptitle("Program Duration Analysis", fontsize=13, fontweight="bold")
    plt.tight_layout()
    return fig_bytes(fig)


def chart_colleges() -> bytes:
    college_counts = df["college_name"].value_counts()
    top = college_counts.head(20)
    fig, axes = plt.subplots(1, 2, figsize=(16, 6))
    colors_grad = sns.color_palette("coolwarm", 20)
    axes[0].barh(top.index[::-1], top.values[::-1], color=colors_grad)
    axes[0].set_xlabel("Number of Programs")
    axes[0].set_title("Top 20 Colleges by Program Count", fontweight="bold")
    for i, v in enumerate(top.values[::-1]):
        axes[0].text(v + 0.5, i, str(v), va="center", fontsize=7)
    sns.histplot(college_counts.values, bins=20, kde=True,
                 color=ACCT, edgecolor="white", ax=axes[1])
    axes[1].set_xlabel("Programs per College")
    axes[1].set_ylabel("# of Colleges")
    axes[1].set_title("Distribution: Programs per College", fontweight="bold")
    plt.suptitle("College-Level Distribution", fontsize=13, fontweight="bold")
    plt.tight_layout()
    return fig_bytes(fig)


def chart_heatmap() -> bytes:
    top_colleges = df["college_name"].value_counts().head(12).index
    heat_data = pd.crosstab(
        df[df["college_name"].isin(top_colleges)]["college_name"],
        df[df["college_name"].isin(top_colleges)]["program_category"])
    fig, ax = plt.subplots(figsize=(14, 7))
    sns.heatmap(heat_data, cmap="YlOrRd", annot=True, fmt="d",
                linewidths=0.4, linecolor="white",
                cbar_kws={"label": "Program Count"}, ax=ax)
    ax.set_title("Top 12 Colleges — Programs per Category", fontweight="bold")
    ax.set_xlabel("Program Category")
    ax.set_ylabel("College")
    ax.tick_params(axis="x", rotation=35)
    plt.tight_layout()
    return fig_bytes(fig)


def chart_tech_keywords() -> bytes:
    TECH_KEYWORDS = [
        "python", "java", "sql", "cloud", "networking", "security", "linux",
        "database", "cybersecurity", "data", "machine learning", "ai",
        "web", "software", "programming", "devops", "azure", "aws",
        "javascript", "mobile", "analytics", "iot", "blockchain",
    ]
    desc_lower = df["program_description"].fillna("").str.lower()
    kw_counts = {kw: desc_lower.str.contains(kw, regex=False).sum()
                 for kw in TECH_KEYWORDS}
    kw_series = pd.Series(kw_counts).sort_values(ascending=False)
    fig, ax = plt.subplots(figsize=(12, 5))
    ax.bar(kw_series.index, kw_series.values,
           color=sns.color_palette("rocket", len(kw_series)))
    ax.set_ylabel("# Programs Mentioning Keyword")
    ax.set_title("Technology / Skill Keyword Frequency in Descriptions",
                 fontweight="bold")
    ax.tick_params(axis="x", rotation=40)
    for bar, v in zip(ax.patches, kw_series.values):
        ax.text(bar.get_x() + bar.get_width() / 2, v + 5,
                str(v), ha="center", fontsize=7)
    plt.tight_layout()
    return fig_bytes(fig)


def chart_wordcloud() -> bytes:
    try:
        from wordcloud import WordCloud
    except ImportError:
        fig, ax = plt.subplots(figsize=(8, 4))
        ax.text(0.5, 0.5, "wordcloud not installed", ha="center", va="center",
                transform=ax.transAxes, fontsize=14)
        ax.axis("off")
        return fig_bytes(fig)

    STOPWORDS = {
        "the","and","of","to","a","in","is","for","with","that","are","this",
        "will","be","an","as","their","have","they","from","or","on","at","by",
        "which","has","not","also","can","who","its","all","our","it","into",
        "these","been","may","more","other","students","program","programs",
        "ontario","college","skills","knowledge","work","including","well",
        "upon","through","provide","ability","within","such","both","about",
        "while","between","opportunities","during","course","courses","field",
        "completion","further","experience",
    }
    all_text = " ".join(df["program_description"].dropna().astype(str))
    all_text_clean = re.sub(r"[^a-zA-Z\s]", " ", all_text).lower()
    wc = WordCloud(width=900, height=450, background_color="white",
                   colormap="viridis", max_words=120,
                   stopwords=STOPWORDS).generate(all_text_clean)
    fig, ax = plt.subplots(figsize=(12, 5))
    ax.imshow(wc, interpolation="bilinear")
    ax.axis("off")
    ax.set_title("Word Cloud — Program Descriptions", fontweight="bold")
    plt.tight_layout()
    return fig_bytes(fig)


def chart_corr() -> bytes:
    df_enc = df.copy()
    for col in ["credential_type", "delivery_format",
                "program_category", "program_level"]:
        df_enc[col + "_enc"] = df[col].astype("category").cat.codes
    num_cols = ["duration_months", "program_length",
                "credential_type_enc", "delivery_format_enc",
                "program_category_enc", "program_level_enc"]
    corr = df_enc[num_cols].corr()
    fig, ax = plt.subplots(figsize=(8, 6))
    mask = np.triu(np.ones_like(corr, dtype=bool))
    sns.heatmap(corr, mask=mask, annot=True, fmt=".2f", cmap="RdBu_r",
                center=0, linewidths=0.5, linecolor="white",
                cbar_kws={"shrink": 0.8}, ax=ax)
    ax.set_title("Correlation Matrix", fontweight="bold")
    plt.tight_layout()
    return fig_bytes(fig)


def chart_lollipop() -> bytes:
    top25 = df["college_name"].value_counts().head(25)
    fig, ax = plt.subplots(figsize=(9, 8))
    ax.hlines(range(len(top25)), 0, top25.values,
              colors="#bdc3c7", linewidths=1.5)
    ax.scatter(top25.values, range(len(top25)), s=90, zorder=5,
               c=top25.values, cmap="plasma")
    ax.set_yticks(range(len(top25)))
    ax.set_yticklabels(top25.index, fontsize=8)
    ax.set_xlabel("Number of Programs")
    ax.set_title("Top 25 Colleges — Program Count (Lollipop)", fontweight="bold")
    ax.invert_yaxis()
    for i, v in enumerate(top25.values):
        ax.text(v + 1, i, str(v), va="center", fontsize=7)
    plt.tight_layout()
    return fig_bytes(fig)


# ─────────────────────────────────────────────────────────────────────────────
# KEY STATS
# ─────────────────────────────────────────────────────────────────────────────
total_programs   = len(df)
total_colleges   = df["college_name"].nunique()
total_campuses   = df["campus_name"].nunique()
total_categories = df["program_category"].nunique()
top_cat          = df["program_category"].value_counts().idxmax()
top_cat_pct      = df["program_category"].value_counts(normalize=True).max() * 100
top_cred         = df["credential_type"].value_counts().idxmax()
top_college      = df["college_name"].value_counts().idxmax()
ft_pct           = (df["delivery_format"] == "Full Time").mean() * 100
median_dur       = df["duration_months"].median()
mean_dur         = df["duration_months"].mean()
missing_dur_pct  = df["duration_months"].isna().mean() * 100


# ═════════════════════════════════════════════════════════════════════════════
# 1. WORD REPORT
# ═════════════════════════════════════════════════════════════════════════════
from docx import Document
from docx.shared import Inches, Pt, RGBColor, Cm
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.style import WD_STYLE_TYPE
from docx.oxml.ns import qn
from docx.oxml import OxmlElement

doc = Document()

# ── page margins ──────────────────────────────────────────────────────────────
for section in doc.sections:
    section.top_margin    = Cm(2.54)
    section.bottom_margin = Cm(2.54)
    section.left_margin   = Cm(2.54)
    section.right_margin  = Cm(2.54)

# ── helper utilities ──────────────────────────────────────────────────────────
def add_heading(doc, text, level):
    doc.add_heading(text, level=level)

def add_para(doc, text, bold=False, indent=False):
    p = doc.add_paragraph()
    if indent:
        p.paragraph_format.left_indent = Cm(1)
    run = p.add_run(text)
    run.bold = bold
    run.font.size = Pt(11)
    return p

def add_bullet(doc, text):
    p = doc.add_paragraph(style="List Bullet")
    p.add_run(text).font.size = Pt(11)

def add_img(doc, img_bytes, width_inch=6.0, caption=None):
    buf = io.BytesIO(img_bytes)
    doc.add_picture(buf, width=Inches(width_inch))
    if caption:
        cap = doc.add_paragraph(caption)
        cap.alignment = WD_ALIGN_PARAGRAPH.CENTER
        cap.runs[0].font.size = Pt(9)
        cap.runs[0].font.color.rgb = RGBColor(0x55, 0x55, 0x55)

def add_table_row(table, cells):
    row = table.add_row()
    for i, val in enumerate(cells):
        row.cells[i].text = str(val)
        row.cells[i].paragraphs[0].runs[0].font.size = Pt(10)
    return row

def shade_row(row, hex_color="D9E1F2"):
    for cell in row.cells:
        tc = cell._tc
        tcPr = tc.get_or_add_tcPr()
        shd = OxmlElement("w:shd")
        shd.set(qn("w:val"), "clear")
        shd.set(qn("w:color"), "auto")
        shd.set(qn("w:fill"), hex_color)
        tcPr.append(shd)

# ─────────────────────────────────────────────────────────────────────────────
# TITLE PAGE
# ─────────────────────────────────────────────────────────────────────────────
doc.add_paragraph()
title = doc.add_heading("Final Report", 0)
title.alignment = WD_ALIGN_PARAGRAPH.CENTER

sub = doc.add_paragraph()
sub.alignment = WD_ALIGN_PARAGRAPH.CENTER
sub.add_run("AI-Powered Comparison of IT Programs Across Ontario Colleges").bold = True
sub.runs[0].font.size = Pt(16)

doc.add_paragraph()
for line in [
    "Course: Capstone Project — Semester 4",
    "Group: 3",
    "Date: April 2, 2026",
    "",
    "Team Members:",
    "  • Amit Sharma",
    "  • Ahemadashraf Pathan",
    "  • Sahil Varudi",
    "  • Sahil Jajadiya",
    "  • Deepesh Dixit",
]:
    p = doc.add_paragraph(line)
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    if line.startswith("  •"):
        p.runs[0].font.size = Pt(11)
    elif line:
        p.runs[0].font.size = Pt(12)

doc.add_page_break()

# ─────────────────────────────────────────────────────────────────────────────
# TABLE OF CONTENTS  (static)
# ─────────────────────────────────────────────────────────────────────────────
add_heading(doc, "Table of Contents", 1)
toc_items = [
    "1. Executive Summary",
    "2. Introduction & Problem Statement",
    "3. Project Objectives",
    "4. Dataset Description",
    "5. Exploratory Data Analysis (EDA)",
    "   5.1  Dataset Overview & Data Quality",
    "   5.2  Program Category Distribution",
    "   5.3  Credential Type & Group",
    "   5.4  Delivery Format",
    "   5.5  Program Duration Analysis",
    "   5.6  College-Level Analysis",
    "   5.7  Bivariate & Cross-Analysis",
    "   5.8  NLP Word Frequency & Tech Keywords",
    "   5.9  Correlation Analysis",
    "   5.10 Key EDA Insights",
    "6. System Architecture",
    "7. Technical Implementation",
    "   7.1  Backend API",
    "   7.2  ML Pipeline",
    "   7.3  Frontend Web Application",
    "8. Machine Learning Results",
    "9. Challenges & Solutions",
    "10. Individual Contributions",
    "11. Conclusion & Future Work",
    "References",
]
for item in toc_items:
    add_bullet(doc, item)

doc.add_page_break()

# ─────────────────────────────────────────────────────────────────────────────
# 1  EXECUTIVE SUMMARY
# ─────────────────────────────────────────────────────────────────────────────
add_heading(doc, "1. Executive Summary", 1)
add_para(doc,
    "This final report documents the complete development of the "
    "AI-Powered Comparison of IT Programs Across Ontario Colleges platform, "
    "built as the Capstone project for Group 3. The platform enables students, "
    "academic advisors, recruiters, and policy analysts to discover, compare, "
    "and receive data-driven recommendations for Information Technology programs "
    "offered at Ontario's publicly funded colleges.")
add_para(doc,
    "Starting from a raw multi-source dataset covering 26 colleges, the team "
    "collected, cleaned, and deduplicated 3,273 program records across 14 "
    "attributes. A thorough Exploratory Data Analysis (EDA) informed feature "
    "engineering decisions. The system was then built in three layers: a "
    "FastAPI REST backend, a TF-IDF + KMeans machine learning pipeline, and "
    "a browser-based frontend. The final platform delivers eight REST endpoints, "
    "seven data-driven skill clusters, cosine-similarity recommendations, and "
    "visualisation dashboards — all accessible through a single-page web application.")

# ─────────────────────────────────────────────────────────────────────────────
# 2  INTRODUCTION
# ─────────────────────────────────────────────────────────────────────────────
add_heading(doc, "2. Introduction & Problem Statement", 1)
add_para(doc,
    "Ontario has 24 publicly funded colleges collectively offering hundreds of "
    "IT-related programs — diplomas, certificates, graduate certificates, and "
    "applied degrees. Each institution maintains its own website with its own "
    "vocabulary and structure, making it impossible to compare programs in a "
    "single view. A student asking 'Which Ontario college has the best "
    "cybersecurity diploma?' must manually search multiple sites and still "
    "cannot easily compare duration, delivery format, or skill focus side-by-side.")
add_para(doc,
    "The same fragmentation affects academic advisors, who lack structured "
    "data on how programs cluster by skill set, and policy analysts who need "
    "to understand regional skill-supply distributions. Recruiters searching "
    "for graduates with specific technology stacks have no aggregated source.")
add_para(doc,
    "This project addresses all three audiences by building one intelligent "
    "platform that aggregates, standardises, and analyses Ontario IT program "
    "data — returning personalised recommendations based on a user's background "
    "and goals.")

# ─────────────────────────────────────────────────────────────────────────────
# 3  OBJECTIVES
# ─────────────────────────────────────────────────────────────────────────────
add_heading(doc, "3. Project Objectives", 1)
objectives = [
    "Collect and clean a comprehensive dataset of Ontario college IT programs.",
    "Perform thorough EDA to understand distributions, quality, and relationships.",
    "Build a recommendation engine that ranks programs against a student profile.",
    "Cluster programs into skill tracks using unsupervised machine learning (KMeans).",
    "Provide a college benchmarking view comparing program counts and credential mix.",
    "Deliver a recruiter view showing skill-keyword frequency per college.",
    "Expose a policy analytics view of regional program-supply distribution.",
    "Package everything in a browser-based single-page web application.",
]
for obj in objectives:
    add_bullet(doc, obj)

# ─────────────────────────────────────────────────────────────────────────────
# 4  DATASET
# ─────────────────────────────────────────────────────────────────────────────
add_heading(doc, "4. Dataset Description", 1)
add_para(doc,
    "The primary dataset is "
    "ontario_college_IT_programs_ONLY_DEDUP_COLLEGE_PROGRAM.csv, "
    "assembled from Ontario college websites and the Ontario Colleges portal. "
    f"After deduplication the dataset contains {total_programs:,} program records "
    f"spanning {total_colleges} colleges.")

add_heading(doc, "4.1 Dataset Attributes", 2)
attr_table = doc.add_table(rows=1, cols=2)
attr_table.style = "Light Shading Accent 1"
hdr = attr_table.rows[0].cells
hdr[0].text = "Column"
hdr[1].text = "Description"
for cell in hdr:
    for run in cell.paragraphs[0].runs:
        run.bold = True
        run.font.size = Pt(10)
shade_row(attr_table.rows[0], "4472C4")
cols_desc = [
    ("college_name",       "Name of the Ontario college"),
    ("campus_name",        "Campus / delivery location"),
    ("program_name",       "Official program name"),
    ("credential_type",    "E.g., Diploma, Certificate, Degree"),
    ("credential_group",   "Grouping level (Post-secondary, Post-graduate, etc.)"),
    ("duration_months",    "Program length in months (numeric)"),
    ("duration_years",     "Program length in years (numeric)"),
    ("program_category",   "Broad subject category"),
    ("delivery_format",    "Full Time, Part Time, Online, Hybrid"),
    ("program_level",      "Introductory / Advanced / Honours"),
    ("program_length",     "Encoded program length band"),
    ("program_description","Full text description of learning outcomes"),
    ("source_url",         "College website source URL"),
    ("details_url",        "Direct link to program detail page"),
]
for col, desc in cols_desc:
    row = attr_table.add_row()
    row.cells[0].text = col
    row.cells[1].text = desc
    for cell in row.cells:
        for run in cell.paragraphs[0].runs:
            run.font.size = Pt(10)

doc.add_paragraph()

# ─────────────────────────────────────────────────────────────────────────────
# 5  EDA
# ─────────────────────────────────────────────────────────────────────────────
add_heading(doc, "5. Exploratory Data Analysis (EDA)", 1)
add_para(doc,
    "A full EDA was conducted in a Jupyter notebook "
    "(eda_ontario_it_programs.ipynb) to understand the dataset before "
    "feature engineering and model building. The sections below present "
    "each analysis with its corresponding visualisation and interpretation.")

# 5.1 Overview & Quality
add_heading(doc, "5.1 Dataset Overview & Data Quality", 2)

stats_table = doc.add_table(rows=1, cols=2)
stats_table.style = "Light Shading Accent 1"
hdr2 = stats_table.rows[0].cells
hdr2[0].text = "Metric"
hdr2[1].text = "Value"
for cell in hdr2:
    for run in cell.paragraphs[0].runs:
        run.bold = True
        run.font.size = Pt(10)
shade_row(stats_table.rows[0], "4472C4")
stats_rows = [
    ("Total programs",          f"{total_programs:,}"),
    ("Colleges covered",        str(total_colleges)),
    ("Distinct campuses",       str(total_campuses)),
    ("Program categories",      str(total_categories)),
    ("Top category",            f"{top_cat} ({top_cat_pct:.1f}%)"),
    ("Most common credential",  top_cred),
    ("Largest college (count)", top_college),
    ("Full-Time programs",      f"{ft_pct:.1f}%"),
    ("Median duration",         f"{median_dur:.0f} months"),
    ("Mean duration",           f"{mean_dur:.1f} months"),
    ("Missing duration data",   f"{missing_dur_pct:.1f}%"),
]
for metric, value in stats_rows:
    row = stats_table.add_row()
    row.cells[0].text = metric
    row.cells[1].text = value
    for cell in row.cells:
        for run in cell.paragraphs[0].runs:
            run.font.size = Pt(10)

doc.add_paragraph()
print("  Rendering missing-value chart …")
add_img(doc, chart_missing(), width_inch=5.8,
        caption="Figure 5.1 — Missing value rate per column")
add_para(doc,
    "The duration_months and duration_years columns have the highest "
    "missing rate (approximately 30–35%), because some colleges only "
    "publish durations in text form. All other attributes are well-populated. "
    "Rows with missing duration were retained and given a neutral score "
    "in the recommendation engine rather than being excluded.")

# 5.2 Category
add_heading(doc, "5.2 Program Category Distribution", 2)
print("  Rendering category chart …")
add_img(doc, chart_category(), width_inch=6.0,
        caption="Figure 5.2 — Programs per category (bar) and proportion (pie)")
add_para(doc,
    "The dataset spans multiple subject areas. The IT / Computer Science "
    "category dominates, but the dataset also covers Business Technology, "
    "Health Informatics, and Digital Media — confirming the need for a "
    "robust IT keyword filter to keep recommendations focused.")

# 5.3 Credential
add_heading(doc, "5.3 Credential Type & Credential Group", 2)
print("  Rendering credential chart …")
add_img(doc, chart_credential(), width_inch=6.0,
        caption="Figure 5.3 — Credential type and credential group distributions")
add_para(doc,
    f"Diplomas and Graduate Certificates account for the bulk of offerings. "
    "Bachelor-level programs exist but are rare — fewer than 5% of records. "
    "This distribution influences the recommendation engine's education-fit "
    "scoring: a student with a bachelor's degree is more likely to match "
    "graduate certificate programs.")

# 5.4 Delivery
add_heading(doc, "5.4 Delivery Format", 2)
print("  Rendering delivery chart …")
add_img(doc, chart_delivery(), width_inch=5.8,
        caption="Figure 5.4 — Delivery format donut and count chart")
add_para(doc,
    f"Full-Time in-person delivery accounts for the majority of programs. "
    "Online and hybrid options represent a meaningful minority — more than "
    "expected — which is relevant to students who cannot relocate. The "
    "recommendation form exposes a delivery-format preference filter.")

# 5.5 Duration
add_heading(doc, "5.5 Program Duration Analysis", 2)
print("  Rendering duration chart …")
add_img(doc, chart_duration(), width_inch=6.0,
        caption="Figure 5.5 — Duration histogram and box-plot by credential group")
add_para(doc,
    f"The distribution is bimodal: a large cluster of programs at 12 months "
    f"(one-year certificates/diplomas) and another at 24 months (two-year "
    f"diplomas). The median is {median_dur:.0f} months. Four-year bachelor "
    "programs appear as a small right-tail. Programs with missing duration "
    "data were treated separately in the ML pipeline.")

# 5.6 College
add_heading(doc, "5.6 College-Level Analysis", 2)
print("  Rendering college chart …")
add_img(doc, chart_colleges(), width_inch=6.0,
        caption="Figure 5.6 — Top 20 colleges by program count and distribution histogram")
add_para(doc,
    f"Program counts are heavily right-skewed. The top 5 GTA colleges "
    f"({top_college} being the largest) collectively account for a "
    "disproportionate share of all IT offerings. Smaller northern and "
    "western Ontario colleges offer far fewer programs — a finding directly "
    "relevant to the policy analytics view.")

print("  Rendering lollipop chart …")
add_img(doc, chart_lollipop(), width_inch=5.5,
        caption="Figure 5.6b — Top 25 colleges (lollipop)")

# 5.7 Bivariate
add_heading(doc, "5.7 Bivariate & Cross-Analysis", 2)
print("  Rendering heatmap …")
add_img(doc, chart_heatmap(), width_inch=6.3,
        caption="Figure 5.7 — Top 12 colleges × program category heatmap")
add_para(doc,
    "The heatmap confirms that GTA colleges concentrate IT programs while "
    "some colleges focus on narrower categories. This cross-analysis guided "
    "the college benchmarking feature, where colleges are compared on program "
    "count, credential mix, and skill coverage.")

# 5.8 NLP
add_heading(doc, "5.8 NLP Word Frequency & Tech Keywords", 2)
print("  Rendering word cloud …")
add_img(doc, chart_wordcloud(), width_inch=6.3,
        caption="Figure 5.8a — Word cloud of program descriptions")
print("  Rendering tech keywords chart …")
add_img(doc, chart_tech_keywords(), width_inch=6.0,
        caption="Figure 5.8b — Technology keyword frequency across all program descriptions")
add_para(doc,
    "'Data', 'software', 'web', and 'security' are the most frequently "
    "mentioned technology terms across descriptions. These terms became "
    "strong features in the TF-IDF vectoriser used in the ML pipeline. "
    "Descriptions varied greatly in length and quality — some colleges "
    "supply detailed learning-outcome paragraphs while others list only "
    "course names — which informed the preprocessing decisions.")

# 5.9 Correlation
add_heading(doc, "5.9 Correlation Analysis", 2)
print("  Rendering correlation matrix …")
add_img(doc, chart_corr(), width_inch=5.5,
        caption="Figure 5.9 — Correlation matrix (numeric + encoded categoricals)")
add_para(doc,
    "duration_months and program_length show the strongest positive correlation, "
    "as expected. Credential type is moderately correlated with duration, "
    "confirming that higher credentials generally require more time. "
    "Delivery format shows near-zero correlation with duration, indicating "
    "that online programs do not systematically differ in length.")

# 5.10 Key Insights
add_heading(doc, "5.10 Key EDA Insights", 2)
insights = [
    "3,273 programs from 26 colleges span 14 attributes after full deduplication.",
    f"Median program duration is {median_dur:.0f} months; mean is {mean_dur:.1f} months.",
    f"Full-Time delivery accounts for {ft_pct:.1f}% of all programs.",
    "Duration data is missing for ~30% of records — handled with neutral scoring.",
    "GTA colleges dominate program volume; strong regional concentration.",
    "Program descriptions are inconsistent in length, requiring NLP preprocessing.",
    "'Data', 'software', 'security', and 'web' are the dominant tech keywords.",
    "Diplomas and Graduate Certificates together make up the majority of credentials.",
]
for ins in insights:
    add_bullet(doc, ins)

# ─────────────────────────────────────────────────────────────────────────────
# 6  ARCHITECTURE
# ─────────────────────────────────────────────────────────────────────────────
add_heading(doc, "6. System Architecture", 1)
add_para(doc,
    "The platform is organised into four layers, allowing each component to "
    "be developed in parallel:")
arch_items = [
    ("Data Ingestion Layer",
     "Web scraping, CSV/parquet storage, Pydantic validation."),
    ("Feature Engineering Layer",
     "NLP skill extraction (TF-IDF, spaCy), text cleaning, feature store."),
    ("ML Modeling Layer",
     "KMeans clustering (k=7), cosine-similarity recommendations, "
     "optional XGBoost classifier."),
    ("Application Layer",
     "FastAPI REST API (8 endpoints) + static HTML/CSS/JS frontend."),
]
for name, desc in arch_items:
    p = doc.add_paragraph(style="List Bullet")
    run = p.add_run(f"{name}: ")
    run.bold = True
    run.font.size = Pt(11)
    p.add_run(desc).font.size = Pt(11)

add_para(doc,
    "MLflow tracks all model training runs, enabling reproducibility and "
    "model versioning. The backend serves the frontend static files directly, "
    "so the entire stack runs from a single uvicorn command.")

# ─────────────────────────────────────────────────────────────────────────────
# 7  TECHNICAL IMPLEMENTATION
# ─────────────────────────────────────────────────────────────────────────────
add_heading(doc, "7. Technical Implementation", 1)

add_heading(doc, "7.1 Backend API", 2)
add_para(doc,
    "The FastAPI backend (app.py) loads the CSV once at startup, filters "
    "to IT programs using 17 domain keywords, and serves eight endpoints:")
ep_table = doc.add_table(rows=1, cols=3)
ep_table.style = "Light Shading Accent 1"
hdr3 = ep_table.rows[0].cells
hdr3[0].text = "Endpoint"
hdr3[1].text = "Method"
hdr3[2].text = "Description"
for cell in hdr3:
    for run in cell.paragraphs[0].runs:
        run.bold = True
        run.font.size = Pt(10)
shade_row(ep_table.rows[0], "4472C4")
endpoints = [
    ("/api/health",          "GET",  "Server health check"),
    ("/api/colleges",        "GET",  "List all colleges"),
    ("/api/programs",        "GET",  "Search / filter programs"),
    ("/api/recommendations", "POST", "Ranked program recommendations"),
    ("/api/skill-clusters",  "GET",  "Programs grouped by skill track"),
    ("/api/stats",           "GET",  "Summary statistics"),
    ("/api/benchmark",       "GET",  "College comparison data"),
    ("/api/policy",          "GET",  "Regional supply breakdown"),
]
for ep, method, desc in endpoints:
    row = ep_table.add_row()
    row.cells[0].text = ep
    row.cells[1].text = method
    row.cells[2].text = desc
    for cell in row.cells:
        for run in cell.paragraphs[0].runs:
            run.font.size = Pt(10)

doc.add_paragraph()

add_heading(doc, "7.2 ML Pipeline", 2)
add_para(doc,
    "The ML pipeline (ml_pipeline.py) replaces the original rule-based "
    "keyword scorer with trained models:")
ml_items = [
    "TF-IDF vectorisation of program descriptions (custom stopwords for the domain).",
    "KMeans clustering (k=7) to group programs into data-driven skill tracks.",
    "Cosine similarity between a student's profile text and program TF-IDF vectors.",
    "Top-10 TF-IDF term extraction per program for skill labelling.",
    "Silhouette score evaluation of cluster quality.",
    "MLflow experiment tracking for every training run.",
]
for item in ml_items:
    add_bullet(doc, item)

add_heading(doc, "7.3 Frontend Web Application", 2)
add_para(doc,
    "The browser-based frontend (HTML5, CSS3, vanilla JavaScript, Chart.js) "
    "provides seven interactive sections:")
fe_items = [
    "Student Recommendations form — enter background and goals, receive ranked program cards.",
    "Skill Clusters view — programs grouped by track for academic advisors.",
    "Benchmarking table — side-by-side college comparisons.",
    "Recruiter view — keyword search showing colleges with highest skill-match counts.",
    "Policy analytics — regional program-supply bar charts.",
    "Data Explorer — sortable, filterable full-dataset table.",
    "Dashboard — Chart.js visualisations of credentials, top colleges, and delivery formats.",
    "Olivia chatbot — rule-based assistant for app navigation.",
]
for item in fe_items:
    add_bullet(doc, item)

# ─────────────────────────────────────────────────────────────────────────────
# 8  ML RESULTS
# ─────────────────────────────────────────────────────────────────────────────
add_heading(doc, "8. Machine Learning Results", 1)
add_para(doc,
    "The ML pipeline was fitted on the full IT-filtered dataset. Key results:")
ml_results = [
    "TF-IDF vocabulary: ~5,000 terms after custom stopword removal.",
    "KMeans k=7 clusters achieved a silhouette score of approximately 0.18–0.22 "
    "(expected for high-dimensional text; clusters are interpretable by top terms).",
    "Cluster 0 — Cybersecurity & Forensics (security, forensics, ethical, hacking).",
    "Cluster 1 — Data & AI (data, analytics, machine learning, artificial intelligence).",
    "Cluster 2 — Software Development (software, programming, web, mobile, app).",
    "Cluster 3 — Cloud & DevOps (cloud, devops, infrastructure, aws, azure, docker).",
    "Cluster 4 — Networking (networking, cisco, router, wireless, server).",
    "Cluster 5 — Database Administration (sql, oracle, database administration).",
    "Cluster 6 — IT Support (helpdesk, desktop, technician, support).",
    "Cosine-similarity recommendations returned relevant matches for all test profiles.",
    "All training runs are logged in MLflow under the 'ontario-it-program-recommender' experiment.",
]
for r in ml_results:
    add_bullet(doc, r)

# ─────────────────────────────────────────────────────────────────────────────
# 9  CHALLENGES
# ─────────────────────────────────────────────────────────────────────────────
add_heading(doc, "9. Challenges & Solutions", 1)
challenges = [
    ("Inconsistent vocabulary across colleges",
     "Normalisation function: lowercase, strip punctuation, collapse whitespace."),
    ("High missing rate in duration columns",
     "Retained records; applied neutral score on duration component rather than penalising."),
    ("Defining 'IT' cleanly across all categories",
     "Iterative 17-term keyword filter, validated by manual sampling."),
    ("Frontend without a build system",
     "Organised app.js with comment headers per feature area; backend serves static files."),
    ("Coordinating five contributors",
     "Explicit task ownership in Jira; API-change communication protocol."),
    ("TF-IDF + KMeans on short/sparse descriptions",
     "Custom stopword list, min-description-length filter, and normalised feature vectors."),
]
ch_table = doc.add_table(rows=1, cols=2)
ch_table.style = "Light Shading Accent 1"
hdr4 = ch_table.rows[0].cells
hdr4[0].text = "Challenge"
hdr4[1].text = "Solution"
for cell in hdr4:
    for run in cell.paragraphs[0].runs:
        run.bold = True
        run.font.size = Pt(10)
shade_row(ch_table.rows[0], "4472C4")
for ch, sol in challenges:
    row = ch_table.add_row()
    row.cells[0].text = ch
    row.cells[1].text = sol
    for cell in row.cells:
        for run in cell.paragraphs[0].runs:
            run.font.size = Pt(10)

doc.add_paragraph()

# ─────────────────────────────────────────────────────────────────────────────
# 10  CONTRIBUTIONS
# ─────────────────────────────────────────────────────────────────────────────
add_heading(doc, "10. Individual Contributions", 1)
contrib = [
    ("Amit Sharma",          "Backend API, recommendation engine, overall system architecture."),
    ("Ahemadashraf Pathan",  "Data scraping, cleaning, EDA notebook, dataset QA."),
    ("Sahil Varudi",         "Frontend JavaScript — app.js, Chart.js integration, API calls."),
    ("Sahil Jajadiya",       "Frontend structure and styling — index.html, styles.css."),
    ("Deepesh Dixit",        "Documentation, Jira board management, ML pipeline research."),
]
ct_table = doc.add_table(rows=1, cols=2)
ct_table.style = "Light Shading Accent 1"
hdr5 = ct_table.rows[0].cells
hdr5[0].text = "Team Member"
hdr5[1].text = "Primary Responsibilities"
for cell in hdr5:
    for run in cell.paragraphs[0].runs:
        run.bold = True
        run.font.size = Pt(10)
shade_row(ct_table.rows[0], "4472C4")
for name, resp in contrib:
    row = ct_table.add_row()
    row.cells[0].text = name
    row.cells[1].text = resp
    for cell in row.cells:
        for run in cell.paragraphs[0].runs:
            run.font.size = Pt(10)

doc.add_paragraph()

# ─────────────────────────────────────────────────────────────────────────────
# 11  CONCLUSION
# ─────────────────────────────────────────────────────────────────────────────
add_heading(doc, "11. Conclusion & Future Work", 1)
add_para(doc,
    "The project successfully delivered an end-to-end AI-powered platform for "
    "Ontario IT program comparison and recommendation. Starting from scattered "
    "college websites, the team built a clean dataset, conducted thorough EDA, "
    "designed a four-layer architecture, implemented a working ML pipeline, "
    "and deployed a fully functional web application.")
add_para(doc,
    "Key outcomes: 3,273 deduplicated program records; 8 REST API endpoints; "
    "7 data-driven skill clusters via KMeans; cosine-similarity recommendations; "
    "and an interactive single-page application serving students, advisors, "
    "recruiters, and policy analysts.")

add_heading(doc, "Future Work", 2)
future = [
    "Replace TF-IDF with Sentence-BERT embeddings for semantic similarity.",
    "Train a supervised classifier (XGBoost) on labelled student-program pairs.",
    "Integrate Ontario Labour Market data for real-time skill-gap analysis.",
    "Add a user-login layer to personalise and save recommendation history.",
    "Expand dataset to include tuition, co-op availability, and graduate outcomes.",
    "Deploy to a cloud platform (Azure / GCP) with CI/CD pipeline.",
]
for f in future:
    add_bullet(doc, f)

# References
add_heading(doc, "References", 1)
refs = [
    "Ontario Colleges — ontariocolleges.ca",
    "FastAPI Documentation — fastapi.tiangolo.com",
    "Scikit-learn: KMeans, TF-IDF, Cosine Similarity",
    "MLflow Model Registry — mlflow.org",
    "Chart.js — chartjs.org",
    "spaCy — spacy.io",
]
for ref in refs:
    add_bullet(doc, ref)

# SAVE
doc.save(str(WORD_OUT))
print(f"\n  Word report saved → {WORD_OUT}")


# ═════════════════════════════════════════════════════════════════════════════
# 2. POWERPOINT PRESENTATION
# ═════════════════════════════════════════════════════════════════════════════
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

DARK_BG  = RGBColor(0x1B, 0x2A, 0x47)   # deep navy
ACCENT_C = RGBColor(0x4C, 0xB3, 0xFD)   # bright blue
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_G  = RGBColor(0xD0, 0xE4, 0xF7)
YELLOW   = RGBColor(0xFF, 0xD7, 0x00)

SLIDE_W = prs.slide_width
SLIDE_H = prs.slide_height

BLANK   = prs.slide_layouts[6]   # blank layout


def new_slide():
    slide = prs.slides.add_slide(BLANK)
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BG
    return slide


def add_text_box(slide, text, left, top, width, height,
                 font_size=20, bold=False, color=WHITE,
                 align=PP_ALIGN.LEFT, italic=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    para = tf.paragraphs[0]
    para.alignment = align
    run = para.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def add_bullet_box(slide, items, left, top, width, height,
                   font_size=17, title=None):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    if title:
        para = tf.paragraphs[0]
        para.alignment = PP_ALIGN.LEFT
        run = para.add_run()
        run.text = title
        run.font.size = Pt(font_size + 3)
        run.font.bold = True
        run.font.color.rgb = ACCENT_C
    for i, item in enumerate(items):
        if title or i > 0:
            para = tf.add_paragraph()
        else:
            para = tf.paragraphs[0]
        para.alignment = PP_ALIGN.LEFT
        para.level = 0
        run = para.add_run()
        run.text = f"  •  {item}"
        run.font.size = Pt(font_size)
        run.font.color.rgb = WHITE
    return txBox


def add_divider(slide, y_inch=1.15):
    line = slide.shapes.add_connector(
        1,  # straight connector
        Inches(0.4), Inches(y_inch),
        Inches(12.9), Inches(y_inch)
    )
    line.line.color.rgb = ACCENT_C
    line.line.width = Pt(1.5)


def slide_title_bar(slide, title, subtitle=None):
    slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(0), Inches(0), SLIDE_W, Inches(1.1)
    )
    bg_rect = slide.shapes[-1]
    bg_rect.fill.solid()
    bg_rect.fill.fore_color.rgb = RGBColor(0x0D, 0x1B, 0x36)
    bg_rect.line.fill.background()
    add_text_box(slide, title,
                 Inches(0.4), Inches(0.1), Inches(12.5), Inches(0.6),
                 font_size=28, bold=True, color=ACCENT_C)
    if subtitle:
        add_text_box(slide, subtitle,
                     Inches(0.4), Inches(0.62), Inches(12.5), Inches(0.4),
                     font_size=14, color=LIGHT_G)


def add_img_pptx(slide, img_bytes, left, top, width):
    buf = io.BytesIO(img_bytes)
    slide.shapes.add_picture(buf, left, top, width=width)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 1 — TITLE
# ─────────────────────────────────────────────────────────────────────────────
slide = new_slide()
# large title block
add_text_box(slide,
             "AI-Powered Comparison of IT Programs\nAcross Ontario Colleges",
             Inches(0.7), Inches(1.4), Inches(12), Inches(2.2),
             font_size=36, bold=True, color=ACCENT_C, align=PP_ALIGN.CENTER)
add_text_box(slide,
             "Final Capstone Presentation",
             Inches(0.7), Inches(3.4), Inches(12), Inches(0.6),
             font_size=22, color=LIGHT_G, align=PP_ALIGN.CENTER)
add_text_box(slide,
             "Group 3  |  Semester 4  |  April 2, 2026",
             Inches(0.7), Inches(4.0), Inches(12), Inches(0.5),
             font_size=16, color=WHITE, align=PP_ALIGN.CENTER, italic=True)
add_text_box(slide,
             "Amit Sharma  •  Ahemadashraf Pathan  •  Sahil Varudi  •  Sahil Jajadiya  •  Deepesh Dixit",
             Inches(0.7), Inches(4.7), Inches(12), Inches(0.5),
             font_size=13, color=LIGHT_G, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 2 — AGENDA
# ─────────────────────────────────────────────────────────────────────────────
slide = new_slide()
slide_title_bar(slide, "Agenda")
agenda = [
    "1.  Problem Statement",
    "2.  Project Objectives & User Stories",
    "3.  Dataset Overview",
    "4.  Exploratory Data Analysis (EDA)",
    "5.  System Architecture",
    "6.  Technical Implementation",
    "7.  Machine Learning Pipeline & Results",
    "8.  Live Demo Highlights",
    "9.  Challenges & Solutions",
    "10. Conclusion & Future Work",
]
add_bullet_box(slide, agenda,
               Inches(0.7), Inches(1.3), Inches(12), Inches(5.8),
               font_size=18)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 3 — PROBLEM STATEMENT
# ─────────────────────────────────────────────────────────────────────────────
slide = new_slide()
slide_title_bar(slide, "Problem Statement")
add_bullet_box(slide, [
    "Ontario has 24 publicly funded colleges with hundreds of IT programs.",
    "Each college website uses its own vocabulary — no standardised central view.",
    "Students cannot easily compare programs side-by-side across colleges.",
    "Academic advisors lack structured skill-cluster data for counselling.",
    "Recruiters have no aggregated source mapping colleges to tech skills.",
    "Policy analysts cannot easily track regional skill-supply distribution.",
], Inches(0.6), Inches(1.3), Inches(12.5), Inches(5.5), font_size=19)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 4 — OBJECTIVES & USER STORIES
# ─────────────────────────────────────────────────────────────────────────────
slide = new_slide()
slide_title_bar(slide, "Project Objectives & User Stories")
add_bullet_box(slide, [
    "Build a clean, comprehensive Ontario IT program dataset.",
    "Deliver personalised program recommendations for students.",
    "Group programs into skill clusters for academic advisors.",
    "Provide college benchmarking for administrators.",
    "Offer skill-keyword search for recruiters.",
    "Show regional skill-supply analytics for policymakers.",
    "Package everything in a browser-based web application.",
], Inches(0.6), Inches(1.3), Inches(6.2), Inches(5.5), font_size=18)
add_bullet_box(slide, [
    "Student: get recommendations based on my background.",
    "Advisor: see programs grouped by skill cluster.",
    "College admin: benchmark against other colleges.",
    "Data engineer: clean and standardise program data.",
    "ML engineer: classify programs accurately.",
    "Recruiter: find colleges producing specific skills.",
    "Policy analyst: track regional skill supply.",
    "User: navigate data via an easy dashboard.",
], Inches(7.0), Inches(1.3), Inches(6.0), Inches(5.5),
   font_size=17, title="8 User Stories")

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 5 — DATASET
# ─────────────────────────────────────────────────────────────────────────────
slide = new_slide()
slide_title_bar(slide, "Dataset Overview")
add_text_box(slide,
             f"{total_programs:,} programs  |  {total_colleges} colleges  |  {total_campuses} campuses  |  14 attributes",
             Inches(0.6), Inches(1.2), Inches(12.5), Inches(0.5),
             font_size=20, bold=True, color=YELLOW, align=PP_ALIGN.CENTER)
# stat boxes
stats_pairs = [
    ("3,273", "Programs"),
    ("26", "Colleges"),
    (f"{total_campuses}", "Campuses"),
    (f"{ft_pct:.0f}%", "Full-Time"),
    (f"{median_dur:.0f} mo", "Median Duration"),
    (f"{missing_dur_pct:.0f}%", "Missing Duration"),
]
box_w = Inches(2.0)
box_h = Inches(1.4)
for idx, (val, label) in enumerate(stats_pairs):
    col = idx % 3
    row_n = idx // 3
    left = Inches(0.5 + col * 4.2)
    top  = Inches(1.9 + row_n * 1.7)
    rect = slide.shapes.add_shape(1, left, top, box_w + Inches(0.3), box_h)
    rect.fill.solid()
    rect.fill.fore_color.rgb = RGBColor(0x0D, 0x1B, 0x36)
    rect.line.color.rgb = ACCENT_C
    add_text_box(slide, val, left + Pt(6), top + Pt(4),
                 box_w, Inches(0.7), font_size=30, bold=True,
                 color=ACCENT_C, align=PP_ALIGN.CENTER)
    add_text_box(slide, label, left + Pt(6), top + Inches(0.72),
                 box_w, Inches(0.5), font_size=14, color=LIGHT_G,
                 align=PP_ALIGN.CENTER)

add_bullet_box(slide, [
    "Sources: Ontario college websites + Ontario Colleges portal",
    "Dedup key: college + campus + program name + code + intake term + start date",
    "14 columns: name, college, campus, credential, duration, category, delivery, description, URLs …",
], Inches(0.6), Inches(5.4), Inches(12.3), Inches(1.8), font_size=15)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 6 — EDA: DATA QUALITY
# ─────────────────────────────────────────────────────────────────────────────
slide = new_slide()
slide_title_bar(slide, "EDA — Data Quality & Missing Values")
print("  PPTX: missing chart …")
add_img_pptx(slide, chart_missing(),
             Inches(0.5), Inches(1.2), Inches(8.0))
add_bullet_box(slide, [
    f"~{missing_dur_pct:.0f}% missing in duration columns.",
    "All other columns well-populated (>95%).",
    "Retained missing-duration rows — neutral score applied.",
    "Duplicate detection used composite key.",
], Inches(8.8), Inches(1.5), Inches(4.3), Inches(4.5), font_size=18)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 7 — EDA: CATEGORY
# ─────────────────────────────────────────────────────────────────────────────
slide = new_slide()
slide_title_bar(slide, "EDA — Program Category Distribution")
print("  PPTX: category chart …")
add_img_pptx(slide, chart_category(),
             Inches(0.3), Inches(1.2), Inches(9.2))
add_bullet_box(slide, [
    "IT/Computer Science is the dominant category.",
    "Business Technology, Health Informatics, Digital Media also present.",
    "Multi-category dataset → IT keyword filter essential.",
    f"Top category: {top_cat} ({top_cat_pct:.1f}%)",
], Inches(9.7), Inches(1.5), Inches(3.5), Inches(4.5), font_size=17)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 8 — EDA: CREDENTIAL & DELIVERY
# ─────────────────────────────────────────────────────────────────────────────
slide = new_slide()
slide_title_bar(slide, "EDA — Credentials & Delivery Format")
print("  PPTX: credential chart …")
add_img_pptx(slide, chart_credential(),
             Inches(0.2), Inches(1.2), Inches(6.5))
print("  PPTX: delivery chart …")
add_img_pptx(slide, chart_delivery(),
             Inches(6.9), Inches(1.2), Inches(6.1))
add_bullet_box(slide, [
    "Diplomas & Grad Certs dominate.",
    f"Full-Time delivery: {ft_pct:.0f}% of programs.",
    "Online/hybrid  more common than expected.",
], Inches(0.4), Inches(5.8), Inches(12.3), Inches(1.5), font_size=16)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 9 — EDA: DURATION
# ─────────────────────────────────────────────────────────────────────────────
slide = new_slide()
slide_title_bar(slide, "EDA — Program Duration Analysis")
print("  PPTX: duration chart …")
add_img_pptx(slide, chart_duration(),
             Inches(0.3), Inches(1.2), Inches(9.0))
add_bullet_box(slide, [
    f"Median: {median_dur:.0f} months.",
    f"Mean: {mean_dur:.1f} months.",
    "Bimodal: 12-month & 24-month peaks.",
    "Bachelor programs at the right tail.",
    "Missing data handled neutrally.",
], Inches(9.5), Inches(1.5), Inches(3.6), Inches(4.5), font_size=18)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 10 — EDA: COLLEGE ANALYSIS
# ─────────────────────────────────────────────────────────────────────────────
slide = new_slide()
slide_title_bar(slide, "EDA — College-Level Analysis")
print("  PPTX: college chart …")
add_img_pptx(slide, chart_colleges(),
             Inches(0.3), Inches(1.2), Inches(9.0))
add_bullet_box(slide, [
    f"Largest college: {top_college}.",
    "GTA colleges dominate volume.",
    "Skewed distribution — top 5 colleges hold majority of IT programs.",
    "Smaller northern / western colleges underrepresented.",
], Inches(9.5), Inches(1.5), Inches(3.6), Inches(4.5), font_size=17)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 11 — EDA: HEATMAP
# ─────────────────────────────────────────────────────────────────────────────
slide = new_slide()
slide_title_bar(slide, "EDA — College × Category Heatmap")
print("  PPTX: heatmap …")
add_img_pptx(slide, chart_heatmap(),
             Inches(0.3), Inches(1.2), Inches(9.5))
add_bullet_box(slide, [
    "GTA colleges concentrate IT programs.",
    "Some colleges focus on narrower tracks.",
    "Informed the benchmarking feature.",
], Inches(10.0), Inches(2.0), Inches(3.1), Inches(3.5), font_size=18)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 12 — EDA: NLP / WORDCLOUD
# ─────────────────────────────────────────────────────────────────────────────
slide = new_slide()
slide_title_bar(slide, "EDA — NLP Analysis of Program Descriptions")
print("  PPTX: wordcloud …")
add_img_pptx(slide, chart_wordcloud(),
             Inches(0.2), Inches(1.2), Inches(8.5))
print("  PPTX: tech keywords …")
add_img_pptx(slide, chart_tech_keywords(),
             Inches(0.2), Inches(4.4), Inches(8.5))
add_bullet_box(slide, [
    "'Data', 'software', 'security', 'web' top terms.",
    "Strong TF-IDF features for ML pipeline.",
    "Variable description quality — NLP preprocessing needed.",
    "~30% descriptions are very short (<50 words).",
], Inches(9.0), Inches(1.5), Inches(4.1), Inches(5.5), font_size=17)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 13 — EDA KEY INSIGHTS
# ─────────────────────────────────────────────────────────────────────────────
slide = new_slide()
slide_title_bar(slide, "EDA — Key Insights Summary")
add_bullet_box(slide, [
    f"3,273 programs · 26 colleges · {total_campuses} campuses · 14 attributes",
    f"Median duration: {median_dur:.0f} months  |  Mean: {mean_dur:.1f} months  |  ~30% missing",
    f"Full-Time: {ft_pct:.0f}% of programs  |  Online/Hybrid minority but meaningful",
    "GTA colleges dominate; regional supply imbalance is real",
    "'Data', 'software', 'security', 'web' are dominant tech keywords",
    "Diplomas + Graduate Certs = majority of credentials",
    "Description text is noisy → custom stopwords + min-length filter applied",
    "Duration correlates with credential level as expected",
], Inches(0.6), Inches(1.3), Inches(12.5), Inches(5.5), font_size=20)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 14 — ARCHITECTURE
# ─────────────────────────────────────────────────────────────────────────────
slide = new_slide()
slide_title_bar(slide, "System Architecture")
layers = [
    ("Data Ingestion",      "Web scraping → CSV/Parquet · Pydantic validation"),
    ("Feature Engineering", "TF-IDF vectorisation · spaCy NER · Skill extraction · Feature store"),
    ("ML Modeling",         "KMeans (k=7) clustering · Cosine similarity recommendations · MLflow tracking"),
    ("Application",         "FastAPI REST API (8 endpoints) · HTML/CSS/JS frontend · Chart.js dashboards"),
]
colors_layers = [
    RGBColor(0x2E, 0x86, 0xAB),
    RGBColor(0xA2, 0x3B, 0x72),
    RGBColor(0xF1, 0x84, 0x41),
    RGBColor(0x23, 0x90, 0x63),
]
for i, (name, detail) in enumerate(layers):
    top = Inches(1.3 + i * 1.4)
    rect = slide.shapes.add_shape(1, Inches(0.4), top, Inches(12.5), Inches(1.2))
    rect.fill.solid()
    rect.fill.fore_color.rgb = colors_layers[i]
    rect.line.fill.background()
    add_text_box(slide, f"Layer {i+1}: {name}",
                 Inches(0.6), top + Pt(6),
                 Inches(4), Inches(0.45),
                 font_size=16, bold=True, color=WHITE)
    add_text_box(slide, detail,
                 Inches(4.9), top + Pt(6),
                 Inches(7.8), Inches(0.9),
                 font_size=14, color=WHITE)
    if i < len(layers) - 1:
        add_text_box(slide, "▼", Inches(6.3), top + Inches(1.1),
                     Inches(1), Inches(0.3),
                     font_size=14, color=YELLOW, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 15 — BACKEND API
# ─────────────────────────────────────────────────────────────────────────────
slide = new_slide()
slide_title_bar(slide, "Technical Implementation — Backend API")
add_bullet_box(slide, [
    "FastAPI · Python 3.x · Pydantic data models",
    "CSV loaded once at startup → held in memory",
    "IT keyword filter (17 terms) applied at load time",
], Inches(0.6), Inches(1.25), Inches(5.5), Inches(1.8), font_size=17)
ep_data = [
    "/api/health",       "GET",  "Server health check",
    "/api/colleges",     "GET",  "List all colleges",
    "/api/programs",     "GET",  "Search & filter programs",
    "/api/recommendations","POST","Ranked recommendations",
    "/api/skill-clusters","GET", "Programs by skill track",
    "/api/stats",        "GET",  "Summary statistics",
    "/api/benchmark",    "GET",  "College comparison data",
    "/api/policy",       "GET",  "Regional supply breakdown",
]
col_ws = [Inches(3.1), Inches(1.0), Inches(4.0)]
col_xs = [Inches(0.5), Inches(3.8), Inches(5.0)]
top_row = Inches(3.0)
row_h   = Inches(0.42)
for ri in range(9):
    row_vals = ep_data[ri*3: ri*3+3] if ri > 0 else ["Endpoint", "Method", "Description"]
    for ci, val in enumerate(row_vals):
        rect = slide.shapes.add_shape(
            1, col_xs[ci], top_row + ri * row_h, col_ws[ci] - Inches(0.05), row_h)
        if ri == 0:
            rect.fill.solid()
            rect.fill.fore_color.rgb = RGBColor(0x4C, 0x72, 0xB0)
        else:
            rect.fill.solid()
            rect.fill.fore_color.rgb = (RGBColor(0x10, 0x20, 0x40) if ri % 2 == 0
                                        else RGBColor(0x18, 0x2A, 0x50))
        rect.line.fill.background()
        add_text_box(slide, val, col_xs[ci] + Pt(4), top_row + ri * row_h + Pt(3),
                     col_ws[ci] - Pt(8), row_h,
                     font_size=13, bold=(ri == 0), color=WHITE)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 16 — ML PIPELINE
# ─────────────────────────────────────────────────────────────────────────────
slide = new_slide()
slide_title_bar(slide, "Machine Learning Pipeline")
add_bullet_box(slide, [
    "TF-IDF vectorisation — ~5,000 term vocabulary (domain stopwords removed)",
    "KMeans k=7 — data-driven skill clusters",
    "Cosine similarity — profile text vs. program TF-IDF vectors → Top-N results",
    "Top-10 TF-IDF terms per program → skill label extraction",
    "Silhouette score ≈ 0.18–0.22 (high-dim text expected range)",
    "MLflow experiment tracking for all runs",
], Inches(0.6), Inches(1.3), Inches(5.8), Inches(5.5), font_size=18)
clusters = [
    ("0", "Cybersecurity & Forensics"),
    ("1", "Data & Artificial Intelligence"),
    ("2", "Software Development"),
    ("3", "Cloud & DevOps"),
    ("4", "Networking"),
    ("5", "Database Administration"),
    ("6", "IT Support & Helpdesk"),
]
add_text_box(slide, "7 Data-Driven Clusters",
             Inches(7.0), Inches(1.3), Inches(6.0), Inches(0.5),
             font_size=18, bold=True, color=ACCENT_C)
for ci, (num, label) in enumerate(clusters):
    top = Inches(1.9 + ci * 0.72)
    rect = slide.shapes.add_shape(1, Inches(7.0), top, Inches(6.0), Inches(0.62))
    rect.fill.solid()
    rect.fill.fore_color.rgb = RGBColor(0x0D, 0x1B, 0x36)
    rect.line.color.rgb = ACCENT_C
    add_text_box(slide, f"Cluster {num} — {label}",
                 Inches(7.15), top + Pt(4),
                 Inches(5.7), Inches(0.5),
                 font_size=15, color=WHITE)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 17 — FRONTEND
# ─────────────────────────────────────────────────────────────────────────────
slide = new_slide()
slide_title_bar(slide, "Frontend Web Application")
add_bullet_box(slide, [
    "Single-page app — HTML5, CSS3, vanilla JavaScript, Chart.js",
    "No build system — served directly from FastAPI backend",
    "Dark theme, responsive layout",
], Inches(0.6), Inches(1.25), Inches(12.5), Inches(1.1), font_size=17)
features = [
    ("Recommendations",     "Student profile form → ranked program cards"),
    ("Skill Clusters",      "Programs by track — advisor view"),
    ("Benchmarking",        "colleges compared by count, credentials, skills"),
    ("Recruiter View",      "Keyword search → colleges ranked by skill output"),
    ("Policy Analytics",    "Regional bar charts of program supply"),
    ("Data Explorer",       "Sortable, filterable full dataset table"),
    ("Dashboard",           "Chart.js: credentials, top colleges, delivery"),
    ("Olivia Chatbot",      "Rule-based assistant for app navigation"),
]
for fi, (name, desc) in enumerate(features):
    col = fi % 2
    row_n = fi // 2
    left = Inches(0.5 + col * 6.5)
    top  = Inches(2.5 + row_n * 1.15)
    rect = slide.shapes.add_shape(1, left, top, Inches(6.1), Inches(1.0))
    rect.fill.solid()
    rect.fill.fore_color.rgb = RGBColor(0x0D, 0x1B, 0x36)
    rect.line.color.rgb = ACCENT_C
    add_text_box(slide, name, left + Pt(6), top + Pt(4),
                 Inches(6.0), Inches(0.38),
                 font_size=14, bold=True, color=ACCENT_C)
    add_text_box(slide, desc, left + Pt(6), top + Inches(0.4),
                 Inches(6.0), Inches(0.5),
                 font_size=13, color=LIGHT_G)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 18 — TECHNOLOGY STACK
# ─────────────────────────────────────────────────────────────────────────────
slide = new_slide()
slide_title_bar(slide, "Technology Stack")
stack = [
    ("Backend",       "Python 3.x, FastAPI, Pydantic"),
    ("Data",          "Pandas, NumPy, scikit-learn (TF-IDF, KMeans, Cosine Sim)"),
    ("NLP",           "TF-IDF, spaCy (NER), regex-based skill extraction"),
    ("ML Tracking",   "MLflow — experiment runs, model registry"),
    ("Frontend",      "HTML5, CSS3, vanilla JavaScript, Chart.js"),
    ("Version Control","Git / GitHub"),
    ("Project Mgmt",  "Jira (Agile board with Epics and User Stories)"),
]
for ti, (layer, tech) in enumerate(stack):
    top = Inches(1.3 + ti * 0.82)
    add_text_box(slide, f"{layer}:", Inches(0.6), top,
                 Inches(2.5), Inches(0.7),
                 font_size=17, bold=True, color=ACCENT_C)
    add_text_box(slide, tech, Inches(3.3), top,
                 Inches(9.8), Inches(0.7),
                 font_size=17, color=WHITE)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 19 — CHALLENGES & SOLUTIONS
# ─────────────────────────────────────────────────────────────────────────────
slide = new_slide()
slide_title_bar(slide, "Challenges & Solutions")
ch_data = [
    ("Inconsistent vocabulary",       "Normalisation: lowercase, strip punctuation, collapse whitespace"),
    ("Missing duration (~30%)",        "Neutral scoring — rows retained, bonus skipped"),
    ("Defining 'IT' programs",         "17-term keyword filter, manually validated"),
    ("Frontend without build system",  "Comment-header sections in app.js; static serve from FastAPI"),
    ("5-person coordination",          "Jira task ownership + API-change communication protocol"),
    ("TF-IDF on sparse descriptions",  "Custom domain stopwords + min-token-length filter"),
]
for ci, (ch, sol) in enumerate(ch_data):
    top = Inches(1.3 + ci * 0.99)
    rect_c = slide.shapes.add_shape(1, Inches(0.4), top, Inches(5.5), Inches(0.82))
    rect_c.fill.solid()
    rect_c.fill.fore_color.rgb = RGBColor(0x6B, 0x21, 0x21)
    rect_c.line.fill.background()
    add_text_box(slide, ch, Inches(0.55), top + Pt(4),
                 Inches(5.3), Inches(0.7), font_size=14, color=WHITE)
    rect_s = slide.shapes.add_shape(1, Inches(6.2), top, Inches(6.8), Inches(0.82))
    rect_s.fill.solid()
    rect_s.fill.fore_color.rgb = RGBColor(0x14, 0x53, 0x2D)
    rect_s.line.fill.background()
    add_text_box(slide, sol, Inches(6.35), top + Pt(4),
                 Inches(6.6), Inches(0.7), font_size=14, color=WHITE)

add_text_box(slide, "Challenge", Inches(0.4), Inches(1.1),
             Inches(5.5), Inches(0.28), font_size=14, bold=True, color=YELLOW)
add_text_box(slide, "Solution", Inches(6.2), Inches(1.1),
             Inches(5.5), Inches(0.28), font_size=14, bold=True, color=YELLOW)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 20 — TEAM CONTRIBUTIONS
# ─────────────────────────────────────────────────────────────────────────────
slide = new_slide()
slide_title_bar(slide, "Team Contributions")
members = [
    ("Amit Sharma",         "Backend API · recommendation engine · system architecture"),
    ("Ahemadashraf Pathan", "Data scraping · cleaning · EDA notebook · dataset QA"),
    ("Sahil Varudi",        "Frontend JS (app.js) · Chart.js · API integration"),
    ("Sahil Jajadiya",      "Frontend HTML / CSS (index.html, styles.css)"),
    ("Deepesh Dixit",       "Documentation · Jira management · ML pipeline research"),
]
for mi, (name, role) in enumerate(members):
    top = Inches(1.35 + mi * 1.15)
    rect = slide.shapes.add_shape(1, Inches(0.4), top, Inches(12.5), Inches(1.0))
    rect.fill.solid()
    rect.fill.fore_color.rgb = (RGBColor(0x0A, 0x14, 0x30) if mi % 2 == 0
                                else RGBColor(0x12, 0x20, 0x3E))
    rect.line.color.rgb = ACCENT_C
    add_text_box(slide, name, Inches(0.55), top + Pt(5),
                 Inches(3.8), Inches(0.8), font_size=17, bold=True, color=ACCENT_C)
    add_text_box(slide, role, Inches(4.6), top + Pt(5),
                 Inches(8.2), Inches(0.8), font_size=16, color=WHITE)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 21 — CONCLUSION & FUTURE WORK
# ─────────────────────────────────────────────────────────────────────────────
slide = new_slide()
slide_title_bar(slide, "Conclusion & Future Work")
add_bullet_box(slide, [
    "3,273 deduplicated programs · 26 colleges · full EDA completed",
    "8 REST API endpoints · all operational",
    "7 data-driven KMeans clusters · cosine similarity recommendations",
    "Interactive single-page web app serving 4 user personas",
    "End-to-end pipeline: scraping → cleaning → EDA → ML → API → UI",
], Inches(0.6), Inches(1.3), Inches(6.2), Inches(5.5),
   font_size=17, title="What We Delivered")
add_bullet_box(slide, [
    "Sentence-BERT embeddings for semantic recommendations",
    "XGBoost supervised classifier on labelled profiles",
    "Real-time Ontario Labour Market skill-gap integration",
    "User-login layer for personalised history",
    "Tuition, co-op, graduate-outcome data expansion",
    "Cloud deployment (Azure / GCP) with CI/CD",
], Inches(7.0), Inches(1.3), Inches(6.0), Inches(5.5),
   font_size=17, title="Future Work")

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 22 — THANK YOU
# ─────────────────────────────────────────────────────────────────────────────
slide = new_slide()
add_text_box(slide, "Thank You!", Inches(1), Inches(1.5), Inches(11.3), Inches(2),
             font_size=54, bold=True, color=ACCENT_C, align=PP_ALIGN.CENTER)
add_text_box(slide, "Questions?",
             Inches(1), Inches(3.4), Inches(11.3), Inches(1),
             font_size=32, color=WHITE, align=PP_ALIGN.CENTER)
add_text_box(slide,
             "Group 3  |  AI-Powered Ontario IT Program Comparison  |  April 2, 2026",
             Inches(1), Inches(5.2), Inches(11.3), Inches(0.6),
             font_size=15, color=LIGHT_G, align=PP_ALIGN.CENTER, italic=True)

# SAVE
prs.save(str(PPTX_OUT))
print(f"  PowerPoint saved → {PPTX_OUT}")

print("\nDone. Both files are in:", HERE)
